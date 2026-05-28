"""
资料解析服务
支持 PDF、Word(.docx)、Markdown、TXT、PPT 的文本提取
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class MaterialParser:
    """
    多格式资料解析器
    MVP: 优先支持 PDF, Markdown, TXT
    """

    # 支持的文件格式
    SUPPORTED_TYPES = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".md": "markdown",
        ".txt": "text",
        ".pptx": "pptx",
    }

    def parse(self, file_path: str, file_type: str) -> list[dict]:
        """
        解析文件，返回结构化文本块列表
        
        返回格式:
        [{"text": "解析出的文本", "page": 1, "type": "text"}]
        """
        if file_type == "pdf":
            return self._parse_pdf(file_path)
        elif file_type == "docx":
            return self._parse_docx(file_path)
        elif file_type == "pptx":
            return self._parse_pptx(file_path)
        elif file_type in ("markdown", "text"):
            return self._parse_text(file_path)
        else:
            logger.warning(f"不支持的文件类型: {file_type}")
            return []

    def _parse_pdf(self, file_path: str) -> list[dict]:
        """解析 PDF，按页提取文本"""
        try:
            import fitz  # pymupdf
        except ImportError:
            logger.error("pymupdf 未安装。pip install pymupdf")
            return self._fallback(file_path)

        blocks = []
        try:
            doc = fitz.open(file_path)
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text("text").strip()
                if text:
                    blocks.append({
                        "text": text,
                        "page": page_num,
                        "type": "text",
                    })
            doc.close()
        except Exception as e:
            logger.error(f"PDF解析失败: {e}")
            return self._fallback(file_path)

        logger.info(f"PDF解析完成: {file_path} → {len(blocks)} 页")
        return blocks

    def _parse_docx(self, file_path: str) -> list[dict]:
        """解析 Word 文档"""
        try:
            from docx import Document
        except ImportError:
            logger.error("python-docx 未安装。pip install python-docx")
            return self._fallback(file_path)

        try:
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)

            if not paragraphs:
                return self._fallback(file_path)

            # 按段落合并为块（每~500字一块）
            blocks = self._merge_paragraphs(paragraphs)
            logger.info(f"Word解析完成: {file_path} → {len(blocks)} 块")
            return blocks
        except Exception as e:
            logger.error(f"Word解析失败: {e}")
            return self._fallback(file_path)

    def _parse_pptx(self, file_path: str) -> list[dict]:
        """解析 PPT"""
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx 未安装。pip install python-pptx")
            return self._fallback(file_path)

        try:
            prs = Presentation(file_path)
            blocks = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    blocks.append({
                        "text": "\n".join(texts),
                        "page": slide_num,
                        "type": "slide",
                    })
            logger.info(f"PPT解析完成: {file_path} → {len(blocks)} 页")
            return blocks
        except Exception as e:
            logger.error(f"PPT解析失败: {e}")
            return self._fallback(file_path)

    def _parse_text(self, file_path: str) -> list[dict]:
        """解析纯文本/Markdown"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, "r", encoding="gbk") as f:
                    content = f.read()
            except Exception:
                return self._fallback(file_path)

        if not content.strip():
            return []

        # 按双换行分块
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
        blocks = self._merge_paragraphs(paragraphs)
        logger.info(f"文本解析完成: {file_path} → {len(blocks)} 块")
        return blocks

    def _merge_paragraphs(
        self, paragraphs: list[str], target_chars: int = 800
    ) -> list[dict]:
        """合并短段落为合理大小的块"""
        blocks = []
        current = ""
        page = 1

        for para in paragraphs:
            if len(current) + len(para) > target_chars and current:
                blocks.append({"text": current.strip(), "page": page, "type": "text"})
                current = para
                page += 1
            else:
                current += "\n" + para if current else para

        if current.strip():
            blocks.append({"text": current.strip(), "page": page, "type": "text"})

        return blocks

    def _fallback(self, file_path: str) -> list[dict]:
        """所有解析器失败时的降级方案：读取原始文本"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read(10000)  # 最多读10000字
            return [{"text": text, "page": 1, "type": "fallback"}]
        except Exception:
            return [{"text": f"(无法解析: {Path(file_path).name})", "page": 1, "type": "error"}]


# 全局实例
material_parser = MaterialParser()
